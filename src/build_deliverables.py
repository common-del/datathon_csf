"""Build report.pdf, docs/policy_note.pdf, slides.pptx; fix manifest.yml, claims.json, README.
Run from repo root: python src/build_deliverables.py"""
import os, re, json, copy
HERE = os.path.dirname(os.path.abspath(__file__)); ROOT = os.path.abspath(os.path.join(HERE, ".."))
os.chdir(ROOT)
TEAM = "datathon_csf"

# ============================ 1. MD -> PDF ============================
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib.colors import HexColor
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from PIL import Image as PILImage

INK = HexColor("#123B47"); MUT = HexColor("#5A6B70")
SS = getSampleStyleSheet()
ST = {
 "title": ParagraphStyle("t", parent=SS["Title"], fontName="Helvetica-Bold", fontSize=17, leading=21, textColor=INK, spaceAfter=4, alignment=0),
 "meta":  ParagraphStyle("m", parent=SS["Normal"], fontName="Helvetica", fontSize=8.5, leading=11.5, textColor=MUT, spaceAfter=10),
 "h2":    ParagraphStyle("h2", parent=SS["Heading2"], fontName="Helvetica-Bold", fontSize=12.5, leading=15, textColor=INK, spaceBefore=12, spaceAfter=4),
 "h3":    ParagraphStyle("h3", parent=SS["Heading3"], fontName="Helvetica-Bold", fontSize=10.5, leading=13, textColor=INK, spaceBefore=9, spaceAfter=3),
 "body":  ParagraphStyle("b", parent=SS["Normal"], fontName="Helvetica", fontSize=9.3, leading=12.6, spaceAfter=5, alignment=4),
 "bullet":ParagraphStyle("bu", parent=SS["Normal"], fontName="Helvetica", fontSize=9.3, leading=12.6, spaceAfter=3, leftIndent=14, bulletIndent=4),
 "cell":  ParagraphStyle("c", parent=SS["Normal"], fontName="Helvetica", fontSize=8, leading=10),
 "cellh": ParagraphStyle("ch", parent=SS["Normal"], fontName="Helvetica-Bold", fontSize=8, leading=10),
}
def inline(t):
    t = t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    t = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", t)
    t = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"<i>\1</i>", t)
    t = re.sub(r"`(.+?)`", r"<font face='Courier' size='8'>\1</font>", t)
    return t

def md_to_pdf(md_path, pdf_path, two_page=False):
    lines = open(md_path, encoding="utf-8").read().splitlines()
    story, i, first_h1 = [], 0, True
    S = dict(ST)                      # local copy; two_page = the tight one-pager layout
    if two_page:
        S["title"] = ParagraphStyle("t1", parent=S["title"], fontSize=14.5, leading=17, spaceAfter=2)
        S["meta"]  = ParagraphStyle("m1", parent=S["meta"], fontSize=8, leading=10, spaceAfter=6)
        S["h2"]    = ParagraphStyle("h21", parent=S["h2"], fontSize=10.8, leading=12.6, spaceBefore=7, spaceAfter=2)
        S["bullet"]= ParagraphStyle("bu1", parent=S["bullet"], fontSize=8.6, leading=10.9, spaceAfter=2)
        body = ParagraphStyle("b2", parent=ST["body"], fontSize=8.6, leading=10.9, spaceAfter=3)
    else:
        body = ST["body"]
    while i < len(lines):
        ln = lines[i].rstrip()
        if not ln.strip(): i += 1; continue
        if ln.startswith("# ") and first_h1:
            story.append(Paragraph(inline(ln[2:]), S["title"])); first_h1 = False
            meta = []
            j = i + 1
            while j < len(lines) and lines[j].strip() and not lines[j].startswith("#"):
                meta.append(inline(lines[j].strip())); j += 1
            story.append(Paragraph("<br/>".join(meta), S["meta"])); i = j; continue
        if ln.startswith("## "): story.append(Paragraph(inline(ln[3:]), S["h2"])); i += 1; continue
        if ln.startswith("### "): story.append(Paragraph(inline(ln[4:]), S["h3"])); i += 1; continue
        m = re.match(r"!\[.*?\]\((.+?)\)", ln)
        if m:
            p = m.group(1)
            if os.path.exists(p):
                iw, ih = PILImage.open(p).size
                w = 16.5*cm; h = w*ih/iw
                if h > 9.5*cm: h = 9.5*cm; w = h*iw/ih
                story.append(Image(p, width=w, height=h)); story.append(Spacer(1, 5))
            i += 1; continue
        if ln.startswith("|"):
            rows = []
            while i < len(lines) and lines[i].startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not set("".join(cells)) <= set("-: "): rows.append(cells)
                i += 1
            data = [[Paragraph(inline(c), ST["cellh" if r == 0 else "cell"]) for c in row] for r, row in enumerate(rows)]
            t = Table(data, repeatRows=1)
            t.setStyle(TableStyle([("GRID",(0,0),(-1,-1),0.4,HexColor("#C9D4D6")),
                                   ("BACKGROUND",(0,0),(-1,0),HexColor("#EAF1F1")),
                                   ("VALIGN",(0,0),(-1,-1),"TOP"),
                                   ("LEFTPADDING",(0,0),(-1,-1),4),("RIGHTPADDING",(0,0),(-1,-1),4),
                                   ("TOPPADDING",(0,0),(-1,-1),2.5),("BOTTOMPADDING",(0,0),(-1,-1),2.5)]))
            story.append(t); story.append(Spacer(1, 6)); continue
        if ln.startswith("- "):
            story.append(Paragraph(inline(ln[2:]), S["bullet"], bulletText="-")); i += 1; continue
        num = re.match(r"(\d+)\.\s+(.*)", ln)
        if num:
            story.append(Paragraph(inline(num.group(2)), S["bullet"], bulletText=num.group(1)+".")); i += 1; continue
        story.append(Paragraph(inline(ln), body)); i += 1
    mg = 1.35*cm if two_page else 1.7*cm
    doc = SimpleDocTemplate(pdf_path, pagesize=A4, leftMargin=mg, rightMargin=mg,
                            topMargin=1.2*cm if two_page else 1.5*cm, bottomMargin=1.2*cm if two_page else 1.5*cm,
                            title=os.path.basename(pdf_path), author="Team " + TEAM)
    doc.build(story)
    print("wrote", pdf_path)

# SKIP_DOCS=1 leaves report.pdf and docs/policy_note.pdf untouched. Those two are the
# human-finalised documents committed to the repo; the notebook must never overwrite them.
SKIP_DOCS = bool(os.environ.get("SKIP_DOCS"))
if SKIP_DOCS:
    print("SKIP_DOCS set: leaving report.pdf and docs/policy_note.pdf alone")
else:
    md_to_pdf("docs/report_FINAL.md", "report.pdf")
    md_to_pdf("docs/policy_note_FINAL.md", os.path.join("docs", "policy_note.pdf"), two_page=True)

# ============================ 1b. MD -> DOCX (editable) ============================
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

def _runs(par, text, size=10, bold=False):
    """Render **bold**, *italic* and `code` inside one paragraph."""
    for tok in re.split(r"(\*\*.+?\*\*|(?<!\*)\*[^*]+?\*(?!\*)|`[^`]+?`)", text):
        if not tok: continue
        r = par.add_run()
        if tok.startswith("**") and tok.endswith("**"): r.text, r.bold = tok[2:-2], True
        elif tok.startswith("`") and tok.endswith("`"):
            r.text = tok[1:-1]; r.font.name = "Consolas"; r.font.size = Pt(size - 1.2)
        elif tok.startswith("*") and tok.endswith("*"): r.text, r.italic = tok[1:-1], True
        else: r.text = tok
        if bold: r.bold = True
        if not r.font.size: r.font.size = Pt(size)

def md_to_docx(md_path, docx_path, compact=False):
    """compact=True is the one-page policy-note layout."""
    lines = open(md_path, encoding="utf-8").read().splitlines()
    doc = Document()
    sec = doc.sections[0]
    m = Cm(1.25) if compact else Cm(1.9)
    sec.top_margin = sec.bottom_margin = m
    sec.left_margin = sec.right_margin = Cm(1.4) if compact else Cm(2.0)
    st = doc.styles["Normal"]; st.font.name = "Calibri"
    st.font.size = Pt(9) if compact else Pt(10)
    pf = st.paragraph_format
    if compact:
        pf.space_after = Pt(3); pf.space_before = Pt(0); pf.line_spacing = 1.0
        for nm, sz in (("Heading 1", 12), ("Heading 2", 10.5), ("Title", 15)):
            try:
                h = doc.styles[nm]; h.font.size = Pt(sz)
                h.paragraph_format.space_before = Pt(6); h.paragraph_format.space_after = Pt(2)
            except KeyError:
                pass
    base = 9 if compact else 10
    i, first_h1 = 0, True
    while i < len(lines):
        ln = lines[i].rstrip()
        if not ln.strip(): i += 1; continue
        if ln.startswith("# ") and first_h1:
            h = doc.add_heading(ln[2:], level=0); first_h1 = False
            for r in h.runs: r.font.color.rgb = RGBColor(0x12, 0x3B, 0x47)
            i += 1; continue
        if ln.startswith("## ") or ln.startswith("### "):
            lvl = 1 if ln.startswith("## ") else 2
            h = doc.add_heading(ln[lvl + 2:] if lvl == 2 else ln[3:], level=lvl)
            for r in h.runs: r.font.color.rgb = RGBColor(0x12, 0x3B, 0x47)
            i += 1; continue
        m = re.match(r"!\[.*?\]\((.+?)\)", ln)
        if m:
            p = m.group(1)
            if os.path.exists(p):
                iw, ih = PILImage.open(p).size
                doc.add_picture(p, width=Cm(17.0))
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            i += 1; continue
        if ln.startswith("|"):
            rows = []
            while i < len(lines) and lines[i].startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not set("".join(cells)) <= set("-: "): rows.append(cells)
                i += 1
            if not rows: continue
            ncol = max(len(r) for r in rows)
            t = doc.add_table(rows=0, cols=ncol); t.style = "Table Grid"
            t.alignment = WD_TABLE_ALIGNMENT.CENTER
            for ri, row in enumerate(rows):
                cs = t.add_row().cells
                for ci in range(ncol):
                    cell = cs[ci]; cell.text = ""
                    _runs(cell.paragraphs[0], row[ci] if ci < len(row) else "", size=8.5, bold=(ri == 0))
            doc.add_paragraph(); continue
        if ln.startswith("- "):
            p = doc.add_paragraph(style="List Bullet"); _runs(p, ln[2:], size=base); i += 1; continue
        num = re.match(r"(\d+)\.\s+(.*)", ln)
        if num:
            p = doc.add_paragraph(style="List Number"); _runs(p, num.group(2), size=base); i += 1; continue
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        _runs(p, ln, size=base); i += 1
    doc.core_properties.title = "Where Karnataka's children lose mathematics"
    doc.core_properties.author = "Team " + TEAM
    doc.save(docx_path)
    print("wrote", docx_path)

if not SKIP_DOCS:
    md_to_docx("docs/report_FINAL.md", "report.docx")
    md_to_docx("docs/policy_note_FINAL.md", os.path.join("docs", "policy_note.docx"), compact=True)

# ============================ 2. SLIDES ============================
from pptx import Presentation
from pptx.util import Emu
EMU = 914400
prs = Presentation("docs/slides_TEMPLATE.pptx")

def set_shape_text(sh, new_text):
    """Write lines into existing paragraphs, preserving per-paragraph formatting."""
    tf = sh.text_frame
    lines = new_text.split("\n")
    paras = tf.paragraphs
    for k, para in enumerate(paras):
        txt = lines[k] if k < len(lines) else ""
        if para.runs:
            para.runs[0].text = txt
            for r in para.runs[1:]: r.text = ""
        elif txt:
            para.add_run().text = txt
    for k in range(len(paras), len(lines)):
        para = tf.add_paragraph()
        para.add_run().text = lines[k]
        if paras and paras[0].runs:
            para.runs[0].font.size = paras[0].runs[0].font.size
            para.runs[0].font.name = paras[0].runs[0].font.name

def put_figure(slide, box_idx, img, blank_idxs, caption_idx=None, caption=None):
    shapes = list(slide.shapes)
    box = shapes[box_idx]
    x, y, w, h = box.left, box.top, box.width, box.height
    iw, ih = PILImage.open(img).size
    ar = iw/ih; bar = w/h
    if ar > bar: nw, nh = w, int(w/ar)
    else: nh, nw = h, int(h*ar)
    slide.shapes.add_picture(img, x + (w-nw)//2, y + (h-nh)//2, nw, nh)
    for bi in blank_idxs: set_shape_text(shapes[bi], "")
    if caption_idx is not None: set_shape_text(shapes[caption_idx], caption)

F = "outputs/figures/"
S = prs.slides
# S0 title
set_shape_text(S[0].shapes[3], "The grade 6 collapse is real.\nThe grade 4 recovery is not.")
set_shape_text(S[0].shapes[4], "Move the remedial block to the multiplicative step in grades 5-6, and never publish a rank without coverage.")
set_shape_text(S[0].shapes[5], "Team datathon_csf  ·  Datathon 2026  ·  1 August 2026")
# S1 why it matters
set_shape_text(S[1].shapes[1], "She learns to add. Then we lose her at multiply.")
set_shape_text(S[1].shapes[2], "A child in a Kalaburagi GP gets 42 of 100 questions right; in Udupi, 74. And the older she gets, the worse it goes: grade 6 multiplication fell from 51% to 35% in three years. The ladder breaks at the multiplicative step, in grades 5 and 6, exactly where PARAKH shows government schools losing to private.")
set_shape_text(S[1].shapes[5], "1,379,087")
set_shape_text(S[1].shapes[6], "student records in this census: rural government schools, grades 4-6, three years")
put_figure(S[1], 7, F+"09_year_grade_heatmap.png", [8, 9], 10, "Mean score by grade and year: the gradient inverted.")
set_shape_text(S[1].shapes[12], "")
# S2 audit
set_shape_text(S[2].shapes[4], "1,379,087 student records")
set_shape_text(S[2].shapes[5], "grades 4-6, 3 years, 20 items per paper, official per-file competency map. One row per child; duplicates kept by design (no student ID) and logged.")
set_shape_text(S[2].shapes[12], "3 joining traps defused")
set_shape_text(S[2].shapes[13], "GP ID is the only canonical key (names repeat); GP and cluster do not nest (1,718 GPs span clusters); external joins refused below 60% match. All logged.")
# S3 finding 1: the collapse
set_shape_text(S[3].shapes[1], "The grade 6 collapse survives every selection test we could throw at it")
put_figure(S[3], 2, F+"10_g6_collapse.png", [3, 4], 5, "Built straight from the 9 organiser CSVs by src/figure_g6_collapse.py. Dashed line: the same 2,182 GPs in all three years.")
set_shape_text(S[3].shapes[6], "G6 multiplication 51 → 35, division 55 → 38 (raw)\n\nConstant-GP panel: -19.7pp and -22.0pp, worse than raw\n\nAt zero coverage change: -19.3pp. Stable-coverage districts: -18.7pp\n\nNewly reached GPs score HIGHER than veterans (+4.6pp)\n\nAnd the G4 division 'gain' (+5.6pp) vanishes in the panel (-0.6pp)")
set_shape_text(S[3].shapes[8], "SO WHAT   Treat this as real deterioration, not measurement noise. The remedial block belongs at grades 5-6 multiplication-division, starting this academic year.")
set_shape_text(S[3].shapes[9], "outputs/tables/hypothesis_menu.csv rows EH1-EH6, EH22  ·  g6_collapse_from_raw.csv  ·  items change yearly: competency-level reading, externally corroborated by ASER and PARAKH")
# S4 finding 2: bottleneck
set_shape_text(S[4].shapes[1], "Multiplication gates everything after it")
put_figure(S[4], 2, F+"02_competency_bottleneck.png", [3, 4])
put_figure(S[5], 2, F+"05_floor_vs_mean.png", [3, 4])
put_figure(S[6], 2, F+"04_gender_by_competency.png", [3, 4])
put_figure(S[7], 2, F+"06_bright_spots.png", [3, 4])
put_figure(S[8], 8, F+"07_triage.png", [9, 10])
set_shape_text(S[4].shapes[6], "Bottleneck Score\ngate lift  x  (1 − current mastery)\n\nMultiplication: 55.2% mastery, +45.1pp lift to division\n\nAddition → subtraction lift: +48pp (72% vs 24%)\n\nDivision: weakest safe competency at 43.3%")
set_shape_text(S[4].shapes[8], "SO WHAT   Sequence Kalika Balavardhane's remedial block: fluency first, then multiplication, then division. Grades 5-6 first. Materials by Term 2, 2026-27.")
# S5 finding 3: floor
set_shape_text(S[5].shapes[1], "Averages hide the floor, and the floor tracks who shows up")
set_shape_text(S[5].shapes[12], "12 of 25 districts: floor lagged the mean")
set_shape_text(S[5].shapes[13], "Only Mysuru and Kodagu raised the 10th percentile by 5pp.")
set_shape_text(S[5].shapes[16], "Divergence -6.4pp where coverage grew fastest")
set_shape_text(S[5].shapes[17], "vs +3.0pp where coverage was stable (r = -0.54). A falling floor in an expanding district is triage, not proof of decline.")
set_shape_text(S[5].shapes[19], "SO WHAT   Add two columns to the existing review format: the 10th-percentile score and assessment coverage. Cost: nil. Owner: Samagra Shiksha MIS cell.")
# S6 finding 4: gender
set_shape_text(S[6].shapes[1], "Girls lead on 10 of 11 competencies, and it survives the participation check")
set_shape_text(S[6].shapes[6], "+2.0pp overall; Cohen's d = 0.07\n\nLargest lead: multiplication +3.6pp; only exception: measurement\n\nDistricts with balanced boy-girl coverage: girls +3.6pp\n\nWhere girls are over-assessed the lead SHRINKS (r = -0.47),\nso participation deflates the gap, it does not inflate it")
set_shape_text(S[6].shapes[8], "SO WHAT   Real but small. Protect girls' momentum through the multiplicative wall; do not build a flagship programme on d = 0.07.")
# S7 finding 5: cross-dataset
set_shape_text(S[7].shapes[1], "Independent surveys agree with our map, and some places beat their conditions")
set_shape_text(S[7].shapes[6], "Validation, rural against rural\nASER 2024 district arithmetic: rho = 0.56\nPARAKH RS 2024 govt-school maths: rho = 0.51\n\nByndoor block: +31.5pp above prediction\nBright spots cluster: Belagavi holds 6 of 22 (p = 0.02)")
set_shape_text(S[7].shapes[8], "SO WHAT   Trust the map. Send the CRP cadre to Byndoor and Belagavi's bright blocks with a checklist, then replicate what they find.")
# S8 product
set_shape_text(S[8].shapes[6], "Early-warning model\npersistence features only; CV RMSE 10.5pp, R² 0.49 (1,939 GPs)")
set_shape_text(S[8].shapes[7], "Structural data did not beat the same-as-last-year baseline. We say so, and ship the honest model.")
set_shape_text(S[8].shapes[13], "SO WHAT   One HTML file, no login, no licence: dashboard.html runs offline from the repo, filterable by district and block.")
# S9 recommendations
set_shape_text(S[9].shapes[5], "Rebalance teachers to grade 4-6 PTR, Kalyana Karnataka first")
set_shape_text(S[9].shapes[6], "KK PTR is 35.6 vs 21.8 elsewhere and the gap widened to -15.8pp. Target: PTR under 30 in the 7 districts by June 2028.")
set_shape_text(S[9].shapes[7], "Owner: Commissioner DSEL + KKRDB  ·  2027-28 transfer cycle")
set_shape_text(S[9].shapes[11], "Sequence remediation on multiplication-division")
set_shape_text(S[9].shapes[12], "Grades 5-6, inside Kalika Balavardhane, Term 2 2026-27. Success: G6 multiplication from 35.1% to 45% by the next round.")
set_shape_text(S[9].shapes[13], "Owner: DSERT via BEO + CRP cadre")
set_shape_text(S[9].shapes[18], "Two columns added to the existing format: 10th-percentile score and coverage. Cost: nil.")
set_shape_text(S[9].shapes[19], "Owner: Samagra Shiksha MIS cell  ·  next quarterly review")
# S10 where wrong
set_shape_text(S[10].shapes[7], "Census is 2011, NFHS-5 is 2019-21, assessment is 2022-25. We treat these as slow-moving structural conditions, stated per table.")
set_shape_text(S[10].shapes[12], "Coverage is 37.8%, not a census")
set_shape_text(S[10].shapes[13], "Participation is voluntary and doubled over three years (25.1% to 49.9%). We quantified selection six ways, adjusted every league table for it, and it cannot explain the G6 decline. Denominator reconciles to the cross-validated UDISE file.")
# S11 closer
set_shape_text(S[11].shapes[1], "Put the remedial block at the multiplicative step in grades 5-6, this academic year.")
set_shape_text(S[11].shapes[2], "The Kalaburagi child who can add is still in the room. Reach her before division leaves it.")
set_shape_text(S[11].shapes[3], "Repository: github.com/common-del/datathon_csf   ·   Team datathon_csf   ·   Reproduce with: python src/run_all.py")
prs.save("slides.pptx")
print("wrote slides.pptx")

# ============================ 3. MANIFEST + CLAIMS + README ============================
man = open("manifest.yml", encoding="utf-8").read()
man = man.replace("datathon_csf", TEAM)
man = re.sub(r"(\.\./)+tmp/f3/", "outputs/", man)
man = man.replace("../../../../../tmp/f3/", "outputs/")
# register the standalone raw-CSV replication artefacts, which run_all.py does not produce
extra_outputs = [
 ('  - file: "outputs/figures/10_g6_collapse.png"\n'
  '    description: "headline grade 6 collapse, built directly from the 9 organiser CSVs by src/figure_g6_collapse.py"\n'),
 ('  - file: "outputs/tables/g6_collapse_from_raw.csv"\n'
  '    description: "every number in the headline chart, competency by grade by year, from the raw files"\n'),
 ('  - file: "outputs/tables/g6_collapse_panel_gps.csv"\n'
  '    description: "constant-GP panel: 2,182 GPs present in all three years, grade 6"\n'),
 ('  - file: "outputs/tables/hypothesis_menu.csv"\n'
  '    description: "all 36 hypotheses with verdict, effect, evidence and caveat"\n'),
 ('  - file: "outputs/HYPOTHESIS_REGISTER.xlsx"\n'
  '    description: "hypothesis register: methodology, finding, verdict and assumptions per hypothesis"\n'),
]
for blk in extra_outputs:
    if blk.split('"')[1] not in man:
        man = man.replace("\nexternal_datasets:", "\n" + blk + "\nexternal_datasets:", 1)
open("manifest.yml", "w", encoding="utf-8").write(man)
print("manifest.yml fixed:", "tmp/f3" not in man)

cl = json.load(open("claims.json", encoding="utf-8"))
cl["team_name"] = TEAM
def C(cid, desc, val, out, how):
    return {"claim_id": cid, "description": desc, "claimed_value": val, "unit": "",
            "supporting_output": out, "verification_method": "programmatic",
            "verification_detail": {"how_to_check": how}}
existing = {c["claim_id"] for c in cl["claims"]}
new_claims = [
 C("claim-eh-raw-decline", "G6 multiplicative decline, raw", "multiplication 50.9 to 35.1 (-15.8pp); division 55.0 to 37.6 (-17.4pp), 2022-23 to 2024-25",
   "outputs/tables/state_grade_year.csv", "rows grade=6, years 2022-23 and 2024-25, columns C_multiplication and C_division"),
 C("claim-eh-gradient", "Grade gradient inversion", "2022-23: G4 49.6 < G5 56.2 < G6 62.5; 2024-25: G4 57.8 > G5 51.7 > G6 49.3 (pct_mean)",
   "outputs/tables/state_grade_year.csv", "pct_mean by grade and year"),
 C("claim-district-spread", "District spread, pooled", "Udupi 74.3 to Kalaburagi 42.0 pct_mean",
   "outputs/tables/unit_district_with_external.csv", "pct_mean column, max and min districts"),
 C("claim-eh1-panel", "G6 decline in constant-GP panel", "2,182 GPs in all 3 years: mult 50.5 to 30.8 (-19.7pp), div 54.5 to 32.5 (-22.0pp)",
   "outputs/tables/hypothesis_menu.csv", "row id=EH1, effect field; reproduce via src/extra_hypotheses.py"),
 C("claim-eh2-g4", "G4 division gain is compositional", "raw +5.6pp becomes -0.6pp in constant-GP panel (2,278 GPs)",
   "outputs/tables/hypothesis_menu.csv", "row id=EH2, effect field"),
 C("claim-eh5-intercept", "G6 decline at zero coverage change", "-19.3pp predicted decline at zero district coverage change (n=25)",
   "outputs/tables/hypothesis_menu.csv", "row id=EH5, effect field"),
 C("claim-eh22-newgp", "Newly reached GPs score higher", "new 2024-25 GPs (655, 16% of G6 students) mean 53.1 vs veteran 48.6",
   "outputs/tables/hypothesis_menu.csv", "row id=EH22, effect field"),
 C("claim-eh4-covgrowth", "Coverage growth does not predict score change", "r=-0.18, p=0.396, n=25 districts",
   "outputs/tables/hypothesis_menu.csv", "row id=EH4, effect field"),
 C("claim-coverage-years", "Coverage by year, rural State Govt basis", "25.1% (2022-23), 39.0% (2023-24), 49.9% (2024-25); overall 37.8% = 1,379,087 of 3,644,154",
   "outputs/tables/coverage_district_grade_year.csv", "basis=rural, sum assessed/sum enrolled by year; denominator reconciles to external_data/udise_rural_stategovt_g46_gender_district.csv (1,244,415 in 2022-23)"),
 C("claim-coverage-gender", "Assessment coverage by gender", "girls 39.7%, boys 35.9%, gap +3.8pp; girls better covered in 26 of 31 districts",
   "outputs/tables/coverage_rural_stategovt_gender.csv", "group by g, sum assessed / sum enrolled"),
 C("claim-eh8-signed", "Differential participation deflates the girls lead", "corr(girls-boys coverage gap, girls-boys score gap) = -0.47 (p=1.1e-05, n=80 district-years)",
   "outputs/tables/hypothesis_menu.csv", "row id=EH8, effect field"),
 C("claim-eh7-league", "Coverage-adjusted league table", "17 of 26 districts move >=3 ranks; Raichur 19 to 12, Hassan 4 to 11, Gadag 15 to 20, Kodagu 13 to 8 (2024-25)",
   "outputs/tables/league_coverage_adjusted.csv", "rank_raw vs rank_adj columns"),
 C("claim-eh9-genderbal", "Girls' lead in balance-covered district-years", "+3.57pp across 25 balanced district-years (vs +2.79pp overall)",
   "outputs/tables/hypothesis_menu.csv", "row id=EH9, effect field"),
 C("claim-eh10-genderstable", "Girls' lead stable across years", "+3.17, +2.55, +2.62pp across the three years",
   "outputs/tables/hypothesis_menu.csv", "row id=EH10, effect field"),
 C("claim-eh11-aser", "ASER 2024 agreement", "Spearman rho=0.56 (p=0.004) vs ASER Std6-8 division, n=25 districts",
   "outputs/tables/hypothesis_menu.csv", "row id=EH11, effect field"),
 C("claim-eh19-kk", "Kalyana Karnataka gap widening", "-10.0pp (2022-23), -11.9pp (2023-24), -15.8pp (2024-25), student-weighted",
   "outputs/tables/hypothesis_menu.csv", "row id=EH19, effect field"),
 C("claim-eh14-ladder", "Prerequisite ladder", "addition to subtraction lift +48pp (72% vs 24%); multiplication to division +45pp",
   "outputs/tables/competency_prerequisite_pairs.csv", "rows addition/subtraction and multiplication/division, lift_pp"),
 C("claim-eh16-floorcov", "Floor fall tracks coverage growth", "r=-0.54 (p=0.005, n=25); divergence -6.4pp in fastest tercile vs +3.0pp slowest",
   "outputs/tables/hypothesis_menu.csv", "rows id=EH16 and EH24, effect fields"),
 C("claim-eh23-coastal", "Top 3 pooled districts are the 3 worst-covered", "Udupi pooled rank 1 on 5.7% coverage, Dakshina Kannada rank 2 on 2.0%, Uttara Kannada rank 3 on 0.8%; all three absent from the 2024-25 round",
   "outputs/tables/hypothesis_menu.csv", "row id=EH23, effect field; coverage from coverage_district_grade_year.csv basis=rural"),
 C("claim-eh13-proxy", "Most informative competency", "division and measurement 0.67, multiplication 0.66, shapes 0.58 mean correlation with the other competencies",
   "outputs/tables/hypothesis_menu.csv", "row id=EH13, effect field"),
 C("claim-eh20-items", "Instrument quality", "0 of 180 item-year-grades have discrimination r below 0.15",
   "outputs/tables/item_analysis.csv", "count rows where discrimination_r < 0.15"),
 C("claim-raw-replication", "Headline replicated from raw CSVs by independent code",
   "src/figure_g6_collapse.py reads the 9 organiser CSVs directly and returns G6 multiplication 50.9 to 35.1 and division 55.0 to 37.6; panel of 2,182 GPs gives 50.5 to 30.8 and 54.5 to 32.5",
   "outputs/tables/g6_collapse_from_raw.csv and outputs/tables/g6_collapse_panel_gps.csv",
   "run python src/figure_g6_collapse.py; it shares no code with run_all.py and prints the same numbers"),
 C("claim-coverage-correction", "Coverage denominator corrected during review",
   "src/coverage.py originally joined from the assessed side, dropping district-grade-years with zero assessed children; that reported 41.8% overall and 26.8/43.8/56.2 by year. Corrected to 37.8% and 25.1/39.0/49.9",
   "outputs/tables/coverage_summary.csv", "run python src/fix_coverage.py; it asserts the denominator equals external_data/udise_rural_stategovt_g46_gender_district.csv and refuses to write otherwise"),
 C("claim-row-count-check", "Raw file row counts sum to the analysis dataset",
   "9 files: 104,015 + 156,940 + 209,867 + 113,164 + 161,979 + 204,008 + 95,523 + 151,938 + 181,653 = 1,379,087",
   "outputs/tables/g6_collapse_from_raw.csv", "sum the n_students column across all 9 rows"),
]
# overwrite by id, do not merely append: a stale claim left in place is worse than a missing one
by_id = {c["claim_id"]: c for c in cl["claims"]}
for c in new_claims:
    by_id[c["claim_id"]] = c
cl["claims"] = list(by_id.values())
json.dump(cl, open("claims.json", "w", encoding="utf-8"), indent=2)
print("claims.json:", len(cl["claims"]), "claims, team", cl["team_name"])

rd = open("README.md", encoding="utf-8").read()
if "## Team" not in rd:
    rd = rd.replace("# Datathon 2026", "# Datathon 2026", 1)
    rd = rd.split("\n")
    rd.insert(1, "\n## Team\n\nTeam **datathon_csf** · Datathon 2026 (Akshara Foundation + ACSEL) · Tracks: Data Insights & Visualization, Predictive Analytics, Policy & Intervention Design\n\nPrimary evidence: `report.pdf` and `docs/policy_note.pdf`. Supporting: `slides.pptx`. Every numeric claim: `claims.json`.\n")
    rd = "\n".join(rd)
    open("README.md", "w", encoding="utf-8").write(rd)
print("README updated")
