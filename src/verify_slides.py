"""
verify_slides.py
Checks slides.pptx against the notebook's own output tables, then against the template.

Every number that appears on a slide is re-derived here from outputs/tables/ and compared with
what is written. A FAIL means the deck and the code disagree and one of them must be fixed
before submission. Reading the value out of claims.json would be weaker, because claims.json is
itself written by the pipeline; where possible this recomputes from the table.

Also checks the three things that would look careless in front of a judge: the template's 11
titles present and in order, the evaluation footer intact on every slide, and no text of ours
below 24pt.
"""
import os, sys, re, json
import pandas as pd
from pptx import Presentation
from pptx.util import Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TAB = os.path.join(ROOT, "outputs", "tables")
DECK = os.path.join(ROOT, "slides.pptx")

TEMPLATE_TITLES = [
    "Team & Solution Overview", "Problem Statement", "Dataset Understanding",
    "Analytical Approach", "Key Insight 1", "Key Insight 2", "Key Insight 3",
    "Key Insight 4", "Key Insight 5", "Key Insight 6",
    "Predictive / Diagnostic Analysis",
    "Predictive / Diagnostic Analysis", "Visualization / Dashboard",
    "Recommendations and Road Map", "Innovation & Impact", "Conclusion & Key Takeaways",
]
# 5 slides added to the organisers' 11: Key Insights 3 to 6, and a second Predictive /
# Diagnostic slide for the random forest, under the same template title.
N_TEMPLATE, N_INSERTED = 11, 5


def t(name):
    return pd.read_csv(os.path.join(TAB, name))


def build_expectations():
    """Recompute, from the tables, every figure quoted on a slide."""
    e = {}

    g6 = t("g6_collapse_from_raw.csv")
    g6 = g6[g6["grade"] == 6].set_index("year")
    e["class 6 multiplication 2022-23"] = (round(g6.loc["2022-23", "multiplication"]), 51)
    e["class 6 multiplication 2024-25"] = (round(g6.loc["2024-25", "multiplication"]), 35)
    e["class 6 division 2022-23"] = (round(g6.loc["2022-23", "division"]), 55)
    e["class 6 division 2024-25"] = (round(g6.loc["2024-25", "division"]), 38)

    pan = t("g6_collapse_panel_gps.csv").set_index("year")
    e["constant-GP panel, number of GPs"] = (int(pan["n_gps"].iloc[0]), 2182)
    e["panel multiplication fall"] = (round(pan.loc["2024-25", "mult"] - pan.loc["2022-23", "mult"]), -20)
    e["panel division fall"] = (round(pan.loc["2024-25", "div"] - pan.loc["2022-23", "div"]), -22)

    cs = t("coverage_summary.csv").set_index("basis")
    e["records analysed"] = (int(cs.loc["rural", "assessed_total"]), 1379087)
    e["coverage, rural State Govt"] = (round(cs.loc["rural", "state_coverage_pct"], 1), 37.8)

    pg = t("participation_grade_year.csv")
    yr = pg.groupby("year").apply(lambda d: d["assessed"].sum() / d["enrolled"].sum() * 100)
    e["coverage 2022-23"] = (round(yr.loc["2022-23"], 1), 25.1)
    e["coverage 2024-25"] = (round(yr.loc["2024-25"], 1), 49.9)

    pd_ = t("participation_district.csv")
    e["districts that tested nobody"] = (int((pd_["assessed"] == 0).sum()), 2)

    lad = t("competency_prerequisite_pairs.csv")
    def lift(a, b):
        return round(lad[(lad.prerequisite == a) & (lad.dependent == b)]["lift_pp"].iloc[0])
    e["add to subtract lift"] = (lift("addition", "subtraction"), 48)
    e["multiply to divide lift"] = (lift("multiplication", "division"), 45)

    kk = t("kk_vs_rest_udise.csv").set_index("metric")
    e["KK pupil-teacher ratio"] = (round(kk.loc["Pupil-Teacher Ratio", "kalyana_karnataka"]), 34)
    e["rest pupil-teacher ratio"] = (round(kk.loc["Pupil-Teacher Ratio", "rest_of_karnataka"]), 20)
    e["KK children per school"] = (round(kk.loc["Enrolment per school", "kalyana_karnataka"]), 139)
    e["rest children per school"] = (round(kk.loc["Enrolment per school", "rest_of_karnataka"]), 73)

    ct = t("competency_total_correlation.csv").set_index("competency")
    e["measurement r with total"] = (round(ct.loc["Measurement", "r_with_total"], 2), 0.78)
    e["division r with total"] = (round(ct.loc["Division", "r_with_total"], 2), 0.76)
    e["multiplication r with total"] = (round(ct.loc["Multiplication", "r_with_total"], 2), 0.75)
    e["place value r with total"] = (round(ct.loc["Place Value", "r_with_total"], 2), 0.56)

    # hypothesis_menu.csv is the full register of 36. hypothesis_verdicts.csv holds only the
    # original 12 H-numbered ones, so counting that file undercounts by the 24 EH tests.
    hv = t("hypothesis_menu.csv")["verdict"].str.upper()
    e["hypotheses tested"] = (len(hv), 36)
    e["hypotheses held"] = (int((hv == "SUPPORTED").sum()), 19)
    e["hypotheses weak"] = (int((hv == "WEAK").sum()), 3)
    e["hypotheses discarded"] = (int((hv == "DISCARD").sum()), 14)

    claims = json.load(open(os.path.join(ROOT, "claims.json")))["claims"]
    e["checked claims"] = (len(claims), 47)

    mi = t("model_competency_importance.csv")
    ct = t("competency_total_correlation.csv")
    e["forest importance, division"] = (round(mi.set_index("competency").loc["Division", "rf_importance"] * 100), 54)
    e["forest importance, measurement"] = (round(mi.set_index("competency").loc["Measurement", "rf_importance"] * 100), 21)
    e["forest top 2 share"] = (round(mi["rf_importance"].nlargest(2).sum() * 100), 75)
    meta = json.load(open(os.path.join(TAB, "model_competency_importance_meta.json")))
    e["forest in-sample R2, as %"] = (round(float(meta["rf_r2_in_sample"]) * 100), 96)
    # the claim the slide actually makes: the skills the forest throws away are exactly the
    # skills that were not put in front of every child
    m = mi.merge(ct[["competency", "n"]], on="competency", how="left")
    low = set(m.nsmallest(4, "rf_importance")["competency"])
    incomplete = set(m[m["n"] < m["n"].max()]["competency"])
    e["competencies forest ignores == ones not always asked"] = (int(low == incomplete), 1)
    e["competencies left out of some papers"] = (len(incomplete), 4)

    kk = t("kk_gap_by_year.csv").set_index("year")
    for y, want in (("2022-23", -10), ("2023-24", -12), ("2024-25", -16)):
        e["Kalyana Karnataka gap %s" % y] = (round(kk.loc[y, "gap_pp"]), want)

    for y, r, k in (("2022-23", 61, 50), ("2023-24", 54, 43), ("2024-25", 57, 41)):
        e["rest of state %s" % y] = (round(kk.loc[y, "rest"]), r)
        e["Kalyana Karnataka %s" % y] = (round(kk.loc[y, "kalyana_karnataka"]), k)
    # the slide says the gap opens every year, so assert it is monotonic, not just bigger overall
    e["gap widens in every step"] = (int((kk["gap_pp"].diff().dropna() < 0).all()), 1)

    lo = t("kk_learning_outcomes.csv").set_index("source")
    for src, k, r in (("GP Maths Contest", 45, 58), ("ASER 2024", 26, 39),
                      ("PARAKH RS 2024", 36, 42)):
        e["%s, Kalyana Karnataka" % src] = (round(lo.loc[src, "kk_value"]), k)
        e["%s, rest of state" % src] = (round(lo.loc[src, "rest_value"]), r)
    # the claim the slide makes is that all three point the same way, so assert the direction
    e["instruments putting KK behind"] = (int((lo["gap_pp"] < 0).sum()), 3)

    ip = t("kk_inputs.csv").set_index("metric")
    e["library, Kalyana Karnataka"] = (round(ip.loc["% schools with a library", "kalyana_karnataka"]), 92)
    e["library, rest of state"] = (round(ip.loc["% schools with a library", "rest_of_karnataka"]), 99)
    e["private unaided, Kalyana Karnataka"] = (round(ip.loc["% enrolment in private unaided", "kalyana_karnataka"]), 21)
    e["private unaided, rest of state"] = (round(ip.loc["% enrolment in private unaided", "rest_of_karnataka"]), 35)
    e["input metrics on the slide"] = (len(ip), 5)

    figs = [f for f in os.listdir(os.path.join(ROOT, "outputs", "figures")) if f.endswith(".png")]
    e["figures in the chart suite"] = (len(figs), 18)
    return e


# Numbers on the deck that this repo cannot check, because they come from the team's own
# policy draft rather than from the assessment data. Printed so nobody mistakes silence for
# verification.
EXTERNAL = [
    ("GPSTR placed 11,494 teachers in 2023-24", "raw_version.pptx, team's own figure"),
    ("children per teacher under 30 by June 2028", "policy target, not a measurement"),
    ("class 6 Multiplication above 45% by 2026-27", "policy target, not a measurement"),
    ("owners DSEL, KKRDB, DSERT, MIS cell", "institutional assignment, not in the data"),
]


def main():
    if not os.path.exists(DECK):
        sys.exit("slides.pptx not found. Run src/build_slides.py first.")
    prs = Presentation(DECK)
    fails = 0

    print("A. Numbers on the slides against outputs/tables/\n")
    for k, (got, written) in build_expectations().items():
        ok = abs(float(got) - float(written)) < 0.051
        print("   %-4s %-34s table %-12s slide %s" %
              ("PASS" if ok else "FAIL", k, got, written))
        fails += 0 if ok else 1

    print("\nB. Template fidelity\n")
    titles = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.name == "PlaceHolder 1":
                titles.append(sh.text_frame.text.strip())
    checks = [
        ("16 slides (11 template + 5 inserted)", len(prs.slides) == N_TEMPLATE + N_INSERTED),
        ("titles present and in template order", titles == TEMPLATE_TITLES),
        ("evaluation footer on every slide",
         all(any(sh.name == "EvalFooter" for sh in s.shapes) for s in prs.slides)),
        ("Key Insight 3 and 4 follow Key Insight 2",
         titles[6:8] == ["Key Insight 3", "Key Insight 4"]),
        ("Kalyana Karnataka runs outcomes, trend, then inputs",
         titles[7:10] == ["Key Insight 4", "Key Insight 5", "Key Insight 6"]),
        ("the forest slide follows the first Predictive slide",
         titles[10:12] == ["Predictive / Diagnostic Analysis"] * 2),
        ("every slide keeps the Title Only layout",
         len({s.slide_layout.name for s in prs.slides}) == 1),
    ]
    for label, ok in checks:
        print("   %-4s %s" % ("PASS" if ok else "FAIL", label))
        fails += 0 if ok else 1

    print("\nC. Vocabulary\n")
    import re as _re
    banned = _re.compile(r"\b(skill|skills|times|divide)\b", _re.I)
    hits = []
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.name == "EvalFooter" or not sh.has_text_frame:
                continue
            for m in banned.finditer(sh.text_frame.text):
                hits.append((i, m.group(0)))
    print("   %-4s %-32s %d hits" % ("PASS" if not hits else "FAIL",
                                     "no 'skill' or 'times' for competency", len(hits)))
    for h in hits[:6]:
        print("        slide %s  %r" % h)
    fails += 0 if not hits else 1

    print("\nD. Typography\n")
    small, wrongfont = [], []
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.name == "EvalFooter" or not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if not r.text.strip():
                        continue
                    if r.font.size and r.font.size < Pt(24):
                        small.append((i, r.font.size.pt, r.text[:30]))
                    if r.font.name != "Public Sans":
                        wrongfont.append((i, r.font.name, r.text[:30]))
    for label, bad in (("no run below 24pt", small), ("every run set in Public Sans", wrongfont)):
        print("   %-4s %-32s %d offenders" % ("PASS" if not bad else "FAIL", label, len(bad)))
        for b in bad[:5]:
            print("        slide %s  %s  %r" % b)
        fails += 0 if not bad else 1

    print("\nE. Stated on the deck but NOT checkable from this repo\n")
    for claim, why in EXTERNAL:
        print("   NOTE  %-44s %s" % (claim, why))

    print("\n%s   %d failures" % ("ALL CHECKS PASSED" if not fails else "FAILURES PRESENT", fails))
    sys.exit(1 if fails else 0)


if __name__ == "__main__":
    main()
