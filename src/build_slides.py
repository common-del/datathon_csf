"""
build_slides.py
Fills the organisers' template.pptx with the Datathon 2026 submission.

Rules this file obeys, in order of priority:
  1. The template is the template. All 11 slide titles survive, in the original order, on the
     original layout, with the organisers' evaluation footer untouched on every slide. Two extra
     slides are inserted (Key Insight 3 and 4) using a byte copy of the Key Insight 2 slide, so
     they are structurally identical to the ones the organisers supplied.
  2. Nothing on a slide is smaller than 24pt, except the organisers' own 10pt evaluation footer,
     which is template furniture and not ours to resize. There is an assert for this at the end.
  3. Every figure comes from outputs/slide_figures/, which is built from the notebook's own
     tables. No number is typed in here that is not also in claims.json.
  4. Plain white, Public Sans, no decoration.

Language target: a 10-year-old should follow it. Short sentences. Digits, not words.
"""
import os, sys, copy, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "outputs", "slide_figures")
REPFIG = os.path.join(ROOT, "outputs", "figures")
TEMPLATE = os.path.join(ROOT, "docs", "slides_TEMPLATE.pptx")
OUTPUT = os.path.join(ROOT, "slides.pptx")

FONT = "Public Sans"
BODY_PT = 24                 # the floor. Nothing of ours goes below this.
TITLE_PT = 40                # template is 44; a few section names need the extra room
INK = RGBColor(0x1A, 0x1A, 0x1A)
LABEL = RGBColor(0x2E, 0x7D, 0x46)

# Left text column, and the figure slot. The slot is exactly the size build_slide_figures.py
# renders at, so the picture is placed 1:1 and a point in the chart is a point on the slide.
# Never fit-to-box here: fitting silently rescales the type and breaks the font floor.
TXT = dict(left=Inches(0.62), top=Inches(1.30), width=Inches(6.15), height=Inches(5.30))
PIC = dict(left=Inches(6.95), top=Inches(1.30), width=Inches(6.15), height=Inches(3.95))
CAP_TOP = 5.38               # caption sits under the figure, above the organisers' footer

# ---------------------------------------------------------------------------
# Slide content. Each entry: template title (verbatim), the template's own bullet
# labels (verbatim), the line under each, and the figure.
# ---------------------------------------------------------------------------
DECK = [
    dict(title="Team & Solution Overview", fig=None, body=[
        ("Team name", "Datathon_CSF"),
        ("Members", "Kapil, Susmit and Aniket"),
        ("Challenge area", "Turn a large school maths test into policy the state can act on"),
        ("One-line solution", "Find the competencies and the places to anchor numeracy work"),
        ("Contact", "Central Square Foundation, New Delhi"),
    ]),

    dict(title="Problem Statement", fig="S2_grade_flip.png", body=[
        ("Educational challenge", "Class 6 children are getting worse at maths, not better."),
        ("Objectives", "Check if the fall is real, find where, and see if the numbers can be "
                       "trusted."),
        ("Key questions", "Did scores really fall? Who is missing from the test? What do we "
                          "fix first?"),
        ("Why it matters", "13,79,087 answer sheets. Karnataka spends real money on what "
                           "these say."),
    ], caption=True),

    dict(title="Dataset Understanding", fig="S3_participation.png", body=[
        ("Dataset summary", "13,79,087 answer sheets from 3 years of the GP Maths Contest."),
        ("Variables", "Class, year, gender, district down to Gram Panchayat, 11 competencies."),
        ("Cleaning", "Only right or wrong kept. A blank means the paper never asked that "
                     "question."),
        ("Assumptions", "Rural State Government schools only. No school ID, so we stop at "
                        "Gram Panchayat."),
    ], caption=True),

    dict(title="Analytical Approach", fig="S9_variation.png", body=[
        ("Workflow", "9 raw files, one clean table, 36 tests, 47 checked numbers."),
        ("Methods", "We compare the same 2,182 Gram Panchayats over 3 years, so new places "
                    "joining cannot fake a trend."),
        ("Models", "Variance split, bright-spot residuals, an early-warning model."),
        ("Validation", "Every number re-checked by the notebook, and against ASER 2024 and "
                       "PARAKH 2024."),
    ], caption=True),

    dict(title="Key Insight 1", fig="S1_g6_collapse.png", body=[
        ("Finding", "Class 6 fell hard. Multiplication 51% to 35%. Division 55% to 38%."),
        ("Evidence", "Inside the same 2,182 Gram Panchayats the fall is bigger, not smaller: "
                     "20 and 22 points. 6 different tests all land between 18 and 22."),
        ("Interpretation", "The fall is real. It is not weaker schools joining the test."),
    ], caption=True),

    dict(title="Key Insight 2", fig="S4_coverage_extremes.png", body=[
        ("Finding", "Only 38% of eligible children were tested. 2 districts tested nobody."),
        ("Evidence", "The 3 best-scoring districts are the 3 worst-tested. Score and testing "
                     "move opposite ways."),
        ("Implications", "Do not rank districts yet. Fix who gets tested first."),
    ], caption=True),

    dict(title="Key Insight 3", inserted=True, fig="S5_prereq_ladder.png", body=[
        ("Finding", "Competencies stack. A child who has Addition gets Subtraction right 48 "
                    "points more often."),
        ("Evidence", "A child who has Multiplication gets Division right 45 points more "
                     "often. Division is the weakest of all, at 43%."),
        ("Implications", "More division practice cannot fix division. Fix Multiplication "
                         "first."),
    ], caption=True),

    dict(title="Key Insight 4", inserted=True, fig="S14_kk_outcomes.png", body=[
        ("Finding", "Three different tests all put Kalyana Karnataka's 7 districts behind "
                    "the rest of the state."),
        ("Evidence", "Our contest 45% against 58%. ASER 26% against 39%. PARAKH 36% "
                     "against 42%."),
        ("Implications", "This is not one odd test. Three different measures, three "
                         "different sets of children, same answer."),
    ], caption=True),

    dict(title="Key Insight 5", inserted=True, fig="S15_kk_widening.png", body=[
        ("Finding", "The gap is not steady. It opens every year: 10 points, then 12, "
                    "then 16."),
        ("Evidence", "The rest of the state dipped and recovered, 61% to 54% to 57%. "
                     "Kalyana Karnataka fell in both years, 50% to 43% to 41%."),
        ("Implications", "Nobody is catching up. Waiting one more round makes the job "
                         "bigger, not smaller."),
    ], caption=True),

    dict(title="Key Insight 6", inserted=True, fig="S8_kk_inputs.png", body=[
        ("Finding", "Those same districts have fewer teachers and more children."),
        ("Evidence", "34 children per teacher against 20. 139 per school against 73. "
                     "Internet in 17% of schools against 39%. Libraries almost match, "
                     "92% against 99%."),
        ("Implications", "Only 21% go to private schools against 35%, so government schools "
                         "carry more children on fewer teachers. Allow for inputs and the "
                         "gap drops to 5."),
    ], caption=True),

    dict(title="Predictive / Diagnostic Analysis", fig="S10_signal_competencies.png", body=[
        ("Model", "We checked how closely each competency moves with a child's whole score."),
        ("Metrics", "Bright-spot model explains 26% of block differences. Early-warning "
                    "model 49%."),
        ("Important predictors", "Measurement 0.78, Division 0.76, Multiplication 0.75. "
                                 "Place Value 0.56."),
        ("Early warning", "1,939 Gram Panchayats scored and flagged for the next round."),
    ], caption=True),

    dict(title="Predictive / Diagnostic Analysis", inserted=True,
         fig="S13_forest_importance.png", body=[
        ("Model", "A computer guessed each child's total from their 11 competency scores."),
        ("Metrics", "It got 96% right. But the total IS those 11 scores added up."),
        ("Important predictors", "It leans on Division (54%) and Measurement (21%)."),
        ("Early warning", "The 4 it ignores were left out of most question papers. Missing "
                          "from the blueprint, not unimportant."),
    ], caption=True),

    dict(title="Visualization / Dashboard", fig="S12_coverage_map.png", body=[
        ("Dashboard screenshots", "18 figures, every one built by one script, so every number "
                                  "traces back."),
        ("KPIs", "38% tested. Class 6 Multiplication 35%. Kalyana Karnataka 16 points behind."),
        ("Maps / charts", "Map of who got tested. Dark red districts tested almost nobody, so "
                          "their scores mean least."),
    ], caption=True),

    dict(title="Recommendations and Road Map", fig="S6_weakest_competencies.png", body=[
        ("Priority actions", "1. Move teachers into Kalyana Karnataka's 7 districts. "
                             "2. Class 5 and 6 catch-up on Multiplication and Division. "
                             "3. Print coverage beside every district ranking."),
        ("Impact", "Children per teacher under 30 by June 2028. Class 6 Multiplication "
                   "above 45% by 2026-27."),
        ("Feasibility", "Owners: DSEL with KKRDB, then DSERT, then the MIS cell. GPSTR "
                        "placed 11,494 teachers in 2023-24."),
    ], caption=True),

    dict(title="Innovation & Impact", fig="S11_gender.png", body=[
        ("Novelty", "We tried 6 ways to disprove our own headline. We publish the 14 ideas "
                    "that failed."),
        ("Scalability", "One notebook. Run All rebuilds every number and chart."),
        ("Limitations", "No school ID, so we cannot see inside a Gram Panchayat. We say "
                        "consistent with, never proves."),
        ("Expected impact", "Help aimed at the 2 competencies that carry the most, where it "
                            "is needed most."),
    ], caption=True),

    dict(title="Conclusion & Key Takeaways", fig=None, body=[
        ("Top takeaways", "1. Class 6 maths is falling, and the fall is real. 2. Almost half "
                          "the children are still missing from the test. 3. Multiplication "
                          "and Division are the weakest and the most important."),
        ("Next steps", "Move teachers. Run the catch-up. Test everyone. Re-check before "
                       "calling a recovery real."),
        ("Thank you", "Every number is in claims.json. "
                      "github.com/common-del/datathon_csf"),
    ]),
]

# Slides flagged inserted=True above are cloned from this template slide. It is a plain Title
# Only slide with the title placeholder, the bullet box and the footer, and nothing else, so a
# clone carries no decoration that does not belong on the slide it becomes.
CLONE_FROM = 5


# ---------------------------------------------------------------------------
def duplicate_slide(prs, index):
    """python-pptx has no copy-slide API, so copy the shape tree by hand.

    Copying an existing Key Insight slide rather than building one from a layout keeps the
    new slides byte-identical in structure to the ones the organisers shipped: same
    placeholder ids, same footer, same geometry."""
    src = prs.slides[index]
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):                       # drop the layout's empty placeholders
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new


def move_slide(prs, old_i, new_i):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[old_i])
    lst.insert(new_i, ids[old_i])


def shape_named(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_title(slide, text):
    sh = shape_named(slide, "PlaceHolder 1")
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for r in list(p.runs)[1:]:
        r._r.getparent().remove(r._r)
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(TITLE_PT)
    r.font.name = FONT
    r.font.bold = True
    r.font.color.rgb = INK
    p.alignment = PP_ALIGN.CENTER


def set_body(slide, rows):
    """Rewrite the template's bullet box. Each row is (template label, our line)."""
    sh = shape_named(slide, "TextBox 4")
    sh.left, sh.top, sh.width, sh.height = (TXT["left"], TXT["top"],
                                            TXT["width"], TXT["height"])
    tf = sh.text_frame
    tf.word_wrap = True
    # the template box is wrap="none" with autofit, which would let a long line run off the
    # slide instead of wrapping inside the 6.15in column
    bp = tf._txBody.bodyPr
    bp.set("wrap", "square")
    for tag in ("spAutoFit", "normAutofit"):
        for el in bp.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag):
            bp.remove(el)

    tf.clear()
    for i, (label, line) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r1 = p.add_run(); r1.text = "%s: " % label
        r1.font.bold = True; r1.font.color.rgb = LABEL
        r2 = p.add_run(); r2.text = line
        r2.font.bold = False; r2.font.color.rgb = INK
        for r in (r1, r2):
            r.font.size = Pt(BODY_PT); r.font.name = FONT


def add_caption(slide, text, top_in):
    box = slide.shapes.add_textbox(PIC["left"], Inches(top_in), PIC["width"], Inches(0.90))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(BODY_PT); r.font.name = FONT
    r.font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)


def place_figure(slide, path, caption=None):
    """Place at exactly the rendered size. No fitting, no scaling."""
    from PIL import Image
    iw, ih = Image.open(path).size
    want = PIC["width"] / PIC["height"]
    if abs(iw / ih - want) > 0.01:
        sys.exit("%s has aspect %.3f, the slot is %.3f. Rebuild the figures."
                 % (os.path.basename(path), iw / ih, want))
    slide.shapes.add_picture(path, PIC["left"], PIC["top"],
                             width=PIC["width"], height=PIC["height"])
    if caption:
        add_caption(slide, caption, CAP_TOP)


FOOTER_TOP = Inches(6.75)


CHARS_PER_LINE = 40
LINE_IN = 24 * 1.2 / 72.0
PARA_GAP_IN = 10 / 72.0


def budget(rows, label):
    lines = sum(max(1, -(-len("%s: %s" % r) // CHARS_PER_LINE)) for r in rows)
    need = lines * LINE_IN + len(rows) * PARA_GAP_IN
    have = TXT["height"] / 914400
    if need > have:
        print("   TEXT OVERFLOW %-32s needs %.2fin, box is %.2fin (%d lines)"
              % (label, need, have, lines))
        return 1
    return 0


def check(prs):
    """Two things that would embarrass us in front of a judge, both cheap to assert.

    The 10pt EvalFooter is the organisers' own text and is exempt: it is template furniture,
    not our content, and resizing it would be a bigger deviation than leaving it."""
    bad = 0
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.name == "EvalFooter":
                continue
            if sh.top is not None and sh.height is not None and sh.top + sh.height > FOOTER_TOP:
                print("   OVERFLOW slide %d: %s ends at %.2fin" %
                      (i, sh.name, (sh.top + sh.height) / 914400)); bad += 1
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() and r.font.size and r.font.size < Pt(24):
                        print("   UNDER 24pt slide %d: %.0fpt %r" %
                              (i, r.font.size.pt, r.text[:40])); bad += 1
    if bad:
        sys.exit("%d layout problems, nothing written" % bad)
    print("\n  check: 0 shapes past the footer, 0 runs under 24pt (footer exempt)")


def main():
    if not os.path.exists(TEMPLATE):
        sys.exit("template not found: %s" % TEMPLATE)
    caps = json.load(open(os.path.join(FIG, "captions.json")))["captions"]
    prs = Presentation(TEMPLATE)

    if len(prs.slides) != 11:
        sys.exit("expected the 11-slide template, found %d" % len(prs.slides))

    for pos in [i for i, d in enumerate(DECK) if d.get("inserted")]:
        duplicate_slide(prs, CLONE_FROM)
        move_slide(prs, len(prs.slides) - 1, pos)

    if len(prs.slides) != len(DECK):
        sys.exit("deck has %d slides, content has %d" % (len(prs.slides), len(DECK)))

    over = sum(budget(spec["body"], spec["title"]) for spec in DECK)
    if over:
        sys.exit("%d slides have more text than the template box holds" % over)

    for slide, spec in zip(prs.slides, DECK):
        set_title(slide, spec["title"])
        set_body(slide, spec["body"])
        fig = spec.get("fig")
        if fig:
            path = fig if os.path.isabs(fig) else os.path.join(FIG, fig)
            cap = spec.get("caption")
            if cap is True:
                cap = caps.get(os.path.basename(path))
            place_figure(slide, path, cap)
        print("  slide %-2d %-34s %s" % (list(prs.slides).index(slide) + 1, spec["title"],
                                         os.path.basename(fig) if fig else "(no figure)"))

    check(prs)
    prs.save(OUTPUT)
    print("\nwrote %s" % os.path.relpath(OUTPUT, ROOT))


if __name__ == "__main__":
    main()
